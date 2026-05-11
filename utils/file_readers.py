import os
import re
import json
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd

CACHE_FILE = os.path.join("data", ".agent_cache.json")

def load_cache():
    if os.path.exists(CACHE_FILE):
        try:
            with open(CACHE_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except:
            return {}
    return {}

def save_cache(cache):
    os.makedirs(os.path.dirname(CACHE_FILE), exist_ok=True)
    try:
        with open(CACHE_FILE, "w", encoding="utf-8") as f:
            json.dump(cache, f)
    except Exception as e:
        print(f"Error saving cache: {e}")

_FILE_CACHE = load_cache()

def read_pdf(file_path):
    text = ""
    try:
        reader = PdfReader(file_path)
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
    except Exception as e:
        print(f"Error al leer PDF {file_path}: {e}")
    return text

def read_docx(file_path):
    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"Error al leer DOCX {file_path}: {e}")
    return text

def read_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error al leer PPTX {file_path}: {e}")
    return text

def read_csv(file_path):
    try:
        df = pd.read_csv(file_path)
        return df.to_string()
    except Exception as e:
        print(f"Error al leer CSV {file_path}: {e}")
        return ""

def read_excel(file_path):
    try:
        df = pd.read_excel(file_path)
        return df.to_string()
    except Exception as e:
        print(f"Error al leer Excel {file_path}: {e}")
        return ""

def read_txt(file_path):
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        print(f"Error al leer TXT {file_path}: {e}")
        return ""



def read_file(file_path, progress_callback=None):
    try:
        mtime = os.path.getmtime(file_path)
    except OSError:
        return ""
        
    # Check cache first
    if file_path in _FILE_CACHE:
        cached_mtime, cached_content = _FILE_CACHE[file_path]
        if cached_mtime == mtime:
            if progress_callback:
                progress_callback(f"⚡ Memoria: `{os.path.basename(file_path)}`")
            return cached_content

    if progress_callback:
        progress_callback(f"📄 Procesando: `{os.path.basename(file_path)}`")

    ext = os.path.splitext(file_path)[1].lower()
    
    # Skip unsupported binary formats (images, etc.)
    if ext in ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.ico', '.svg', '.webp', '.mp3', '.mp4', '.avi', '.mov', '.zip', '.rar', '.exe', '.dll'):
        return ""
    
    content = ""
    if ext == ".pdf":
        content = read_pdf(file_path)
    elif ext == ".docx" or ext == ".doc":
        content = read_docx(file_path)
    elif ext == ".pptx" or ext == ".ppt":
        content = read_pptx(file_path)
    elif ext == ".csv":
        content = read_csv(file_path)
    elif ext == ".xlsx" or ext == ".xls":
        content = read_excel(file_path)
    elif ext == ".txt":
        content = read_txt(file_path)
        
    # Update cache
    _FILE_CACHE[file_path] = [mtime, content]
    save_cache(_FILE_CACHE)
    return content

def read_all_files_in_folder(folder_path, progress_callback=None, read_files=None):
    """Lee todos los archivos soportados en una carpeta + tablas SQLite relacionadas."""
    if not os.path.exists(folder_path):
        return ""
    
    combined_text = ""
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            file_path = os.path.join(root, file)
            ext = os.path.splitext(file_path)[1].lower()
            is_image = ext in ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.ico', '.svg', '.webp', '.mp3', '.mp4', '.avi', '.mov', '.zip', '.rar', '.exe', '.dll')
            try:
                content = read_file(file_path, progress_callback)
            except Exception:
                content = ""
            if read_files is not None and not is_image:
                read_files.append(file_path)
            if content.strip():
                combined_text += f"\n\n--- INICIO ARCHIVO: {file} ---\n"
                combined_text += content
                combined_text += f"\n--- FIN ARCHIVO: {file} ---\n"
    
    # También leer tablas SQLite que coincidan con el nombre de la carpeta
    try:
        from database import list_tables, read_table
        folder_name = os.path.basename(folder_path.rstrip('/\\'))
        if folder_name:
            prefix = 't_' + re.sub(r'[^a-zA-Z0-9_]', '_', folder_name).lower()
            matching = list_tables(prefix)
            for tbl in matching:
                df = read_table(tbl)
                if not df.empty:
                    combined_text += f"\n\n--- INICIO TABLA SQLite: {tbl} ---\n"
                    combined_text += df.to_string(max_rows=50)
                    combined_text += f"\n--- FIN TABLA SQLite: {tbl} ---\n"
                    if read_files is not None:
                        read_files.append(tbl)
    except Exception as e:
        print(f"Error leyendo tablas SQLite: {e}")
                
    return combined_text
